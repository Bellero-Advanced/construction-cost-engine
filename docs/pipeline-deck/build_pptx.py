"""Generate Construction Cost Engine — Data Pipeline deck as .pptx (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Theme ---
BG       = RGBColor(0x0E, 0x11, 0x16)
PAPER    = RGBColor(0x16, 0x1B, 0x22)
PAPER2   = RGBColor(0x1F, 0x26, 0x30)
INK      = RGBColor(0xF5, 0xF1, 0xE8)
INK2     = RGBColor(0xC7, 0xC2, 0xB3)
INK3     = RGBColor(0x7D, 0x7A, 0x72)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_BR = RGBColor(0xFB, 0xBF, 0x24)
TEAL     = RGBColor(0x14, 0xB8, 0xA6)
TEAL_LT  = RGBColor(0x5E, 0xEA, 0xD4)
RED      = RGBColor(0xEF, 0x44, 0x44)
RED_LT   = RGBColor(0xFC, 0xA5, 0xA5)
LINE     = RGBColor(0x2A, 0x2F, 0x3A)

FONT_BODY = "Inter"
FONT_MONO = "JetBrains Mono"

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid(); rect.fill.fore_color.rgb = BG
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK, font=FONT_BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_tag(slide, x, y, text):
    add_text(slide, x, y, Inches(8), Inches(0.3), text,
             size=10, bold=True, color=AMBER_BR, font=FONT_MONO)


def add_card(slide, x, y, w, h, *, fill=PAPER2, border=LINE, border_pt=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(border_pt)
    sh.shadow.inherit = False
    return sh


def add_pill(slide, x, y, text, *, color=AMBER_BR, bg=RGBColor(0x2A, 0x22, 0x10)):
    w = Inches(0.05 + 0.085 * len(text))
    h = Inches(0.27)
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.5
    sh.fill.solid(); sh.fill.fore_color.rgb = bg
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Emu(40000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_MONO; r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = color
    return x + w


def add_footer(slide, idx, total):
    add_text(slide, Inches(0.45), Inches(7.15), Inches(6), Inches(0.3),
             "construction-cost-engine · pipeline deck", size=9, color=INK3, font=FONT_MONO)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             f"{idx:02d} / {total:02d}", size=9, color=INK3, font=FONT_MONO, align=PP_ALIGN.RIGHT)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW * idx // total, Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_arrow(slide, x, y, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = LINE
    sh.line.fill.background()
    sh.shadow.inherit = False


TOTAL = 10
slides = []


def new_slide():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    return s

# ============================================================
# 1. TITLE
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 01 · INTRO")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(1.2),
         "Construction Cost Engine", size=46, bold=True, color=INK)
add_text(s, Inches(0.6), Inches(1.95), Inches(11.5), Inches(1.4),
         "ระบบรวบรวมราคาวัสดุก่อสร้างไทยแบบ multi-source — ตั้งแต่ Input ของ 10 แหล่งข้อมูล "
         "(รัฐ + ค้าปลีก) ผ่าน acquisition pipeline ลง Cloudflare KV, เปิดเป็น public REST + UI "
         "เปรียบเทียบราคา",
         size=17, color=INK2)

# pills row
px = Inches(0.6)
for label, color, bg in [
    ("10 SOURCES",       AMBER_BR, RGBColor(0x2A,0x22,0x10)),
    ("17 MATERIALS",     TEAL_LT,  RGBColor(0x0E,0x2A,0x26)),
    ("77 PROVINCES",     AMBER_BR, RGBColor(0x2A,0x22,0x10)),
    ("CLOUDFLARE WORKERS", INK2,   RGBColor(0x22,0x26,0x2E)),
    ("NEXT.JS 16",       INK2,    RGBColor(0x22,0x26,0x2E)),
]:
    px = add_pill(s, px, Inches(3.4), label, color=color, bg=bg) + Inches(0.1)

# agenda
add_text(s, Inches(0.6), Inches(4.05), Inches(6), Inches(0.4),
         "วาระของสไลด์", size=16, bold=True, color=INK)
agenda = [
    ("02", "Overview — สถาปัตยกรรม end-to-end"),
    ("03", "Input Layer — 10 sources, 17 materials, 77 provinces"),
    ("04", "Acquisition — 4 กลยุทธ์เก็บราคา"),
    ("05", "Canonicalization — apples-to-apples (Option B)"),
    ("06", "Storage — KV schema + TTL + weekly snapshots"),
    ("07", "API Layer — public REST + admin endpoints"),
    ("08", "Consumer Layer — UI pages ทั้งหมด"),
    ("09", "Problems & Mitigations"),
    ("10", "Roadmap"),
]
for i, (n, t) in enumerate(agenda):
    row_y = Inches(4.5 + i * 0.27)
    add_text(s, Inches(0.7), row_y, Inches(0.5), Inches(0.3), n,
             size=11, bold=True, color=AMBER_BR, font=FONT_MONO)
    add_text(s, Inches(1.15), row_y, Inches(10), Inches(0.3), t,
             size=12.5, color=INK2)

add_footer(s, 1, TOTAL)

# ============================================================
# 2. OVERVIEW
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 02 · OVERVIEW")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
         "System at a Glance", size=34, bold=True, color=INK)

# flow boxes
stages = [
    ("INPUT",     "10 Sources",  "3 Government\n7 Modern Trade"),
    ("ACQUIRE",   "4 Strategies","PDF · JSON API\nScrapingBee · Manual"),
    ("NORMALIZE", "Canonicalize","brand/size/grade\nfilter (Option B)"),
    ("STORE",     "Cloudflare KV","30-day TTL\nweekly snapshot"),
    ("SERVE",     "REST + UI",   "/api/prices\n/compare-sources"),
]
box_w = Inches(2.15); box_h = Inches(1.65); gap = Inches(0.15)
start_x = Inches(0.6)
top_y = Inches(2.0)
for i, (tag, title, body) in enumerate(stages):
    x = start_x + (box_w + gap) * i
    add_card(s, x, top_y, box_w, box_h)
    add_pill(s, x + Inches(0.18), top_y + Inches(0.15), tag, color=AMBER_BR, bg=RGBColor(0x2A,0x22,0x10))
    add_text(s, x + Inches(0.18), top_y + Inches(0.55), box_w - Inches(0.36), Inches(0.4),
             title, size=14, bold=True, color=INK)
    add_text(s, x + Inches(0.18), top_y + Inches(0.95), box_w - Inches(0.36), Inches(0.7),
             body, size=10, color=INK2, font=FONT_MONO)
    # arrow between boxes (skip last)
    if i < len(stages) - 1:
        ax = x + box_w + Inches(0.01)
        add_arrow(s, ax, top_y + Inches(0.7), gap - Inches(0.02), Inches(0.3))

# 3 why cards
why = [
    ("Why this exists",
     "ราคาวัสดุไทยกระจาย — รัฐ (CGD/TPSO/DIT) + ค้าปลีก (HomePro/SCG/Dohome ฯลฯ). "
     "ต้องรวมเพื่อทำ BoQ ที่อ้างอิงได้", AMBER),
    ("Why Cloudflare Workers",
     "Edge SSR + KV cache + Browser Rendering + Cron triggers ในที่เดียว. "
     "ไม่มี cold start, scale-to-zero, ราคาน้อย", TEAL),
    ("Why TTL = 30 days",
     "ราคาประเมินทางการเปลี่ยน รายไตรมาส. ดึงบ่อยกว่านี้ = "
     "เปลือง credit + ไม่ได้ value", LINE),
]
cw = Inches(4.0); ch = Inches(2.2); cgap = Inches(0.1)
cstart = Inches(0.6); cy = Inches(4.3)
for i, (title, body, accent) in enumerate(why):
    x = cstart + (cw + cgap) * i
    add_card(s, x, cy, cw, ch, border=accent, border_pt=1.5)
    add_text(s, x + Inches(0.25), cy + Inches(0.2), cw - Inches(0.5), Inches(0.4),
             title, size=15, bold=True, color=INK)
    add_text(s, x + Inches(0.25), cy + Inches(0.7), cw - Inches(0.5), ch - Inches(0.8),
             body, size=11.5, color=INK2)
add_footer(s, 2, TOTAL)

# ============================================================
# 3. INPUT
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 03 · INPUT LAYER")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
         "10 Sources × 17 Materials × 77 Provinces", size=30, bold=True, color=INK)

# Two columns
col_w = Inches(6.0); col_h = Inches(3.6)
add_card(s, Inches(0.6), Inches(2.0), col_w, col_h)
add_card(s, Inches(6.8), Inches(2.0), col_w, col_h)

add_text(s, Inches(0.85), Inches(2.15), Inches(5.5), Inches(0.4),
         "🏛 Government (3)", size=16, bold=True, color=AMBER_BR)
gov = [
    ("TPSO (CMI Index)", "PDF",         "Monthly"),
    ("CGD (กรมบัญชีกลาง)", "PDF / blocked", "Quarterly"),
    ("DIT (ก.พาณิชย์)",   "HTML / down", "Daily"),
]
# headers
add_text(s, Inches(0.85), Inches(2.65), Inches(2.2), Inches(0.3), "SOURCE", size=9, bold=True, color=AMBER_BR, font=FONT_MONO)
add_text(s, Inches(3.2),  Inches(2.65), Inches(1.5), Inches(0.3), "FORMAT", size=9, bold=True, color=AMBER_BR, font=FONT_MONO)
add_text(s, Inches(4.9),  Inches(2.65), Inches(1.5), Inches(0.3), "UPDATE", size=9, bold=True, color=AMBER_BR, font=FONT_MONO)
for i, (src, fmt, upd) in enumerate(gov):
    y = Inches(3.05 + i * 0.42)
    add_text(s, Inches(0.85), y, Inches(2.4), Inches(0.4), src, size=11.5, color=INK)
    add_text(s, Inches(3.2),  y, Inches(1.7), Inches(0.4), fmt, size=10.5, color=INK2, font=FONT_MONO)
    add_text(s, Inches(4.9),  y, Inches(1.5), Inches(0.4), upd, size=10.5, color=INK2)
add_text(s, Inches(0.85), Inches(4.7), Inches(5.4), Inches(0.5),
         "→ ราคาอ้างอิงทางการ ใช้ใน BoQ จริง",
         size=11, color=INK3)

add_text(s, Inches(7.05), Inches(2.15), Inches(5.5), Inches(0.4),
         "🛒 Modern Trade (7)", size=16, bold=True, color=TEAL_LT)
retail = [
    ("HomePro",          "JSON suggest.jsp", "LIVE",   TEAL),
    ("MegaHome",         "JSON suggest.jsp", "LIVE",   TEAL),
    ("Thai Watsadu",     "SPA + CF block",   "MANUAL", RED),
    ("Global House",     "SPA + XHR",        "MANUAL", RED),
    ("BnB / SCG / Dohome","SPA",             "MANUAL", RED),
]
add_text(s, Inches(7.05), Inches(2.65), Inches(2.5), Inches(0.3), "SOURCE", size=9, bold=True, color=AMBER_BR, font=FONT_MONO)
add_text(s, Inches(9.55), Inches(2.65), Inches(1.7), Inches(0.3), "FORMAT", size=9, bold=True, color=AMBER_BR, font=FONT_MONO)
add_text(s, Inches(11.3), Inches(2.65), Inches(1.3), Inches(0.3), "STATUS", size=9, bold=True, color=AMBER_BR, font=FONT_MONO)
for i, (src, fmt, st, col) in enumerate(retail):
    y = Inches(3.0 + i * 0.34)
    add_text(s, Inches(7.05), y, Inches(2.5), Inches(0.3), src, size=11, color=INK)
    add_text(s, Inches(9.55), y, Inches(1.8), Inches(0.3), fmt, size=10, color=INK2, font=FONT_MONO)
    bg = RGBColor(0x0E,0x2A,0x26) if col is TEAL else RGBColor(0x2A,0x10,0x10)
    add_pill(s, Inches(11.3), y, st, color=(TEAL_LT if col is TEAL else RED_LT), bg=bg)
add_text(s, Inches(7.05), Inches(4.95), Inches(5.4), Inches(0.5),
         "→ ราคาตลาด สำหรับเทียบ market price",
         size=11, color=INK3)

# Materials taxonomy
add_text(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.4),
         "Materials taxonomy", size=14, bold=True, color=INK)
tax = [
    ("wall_tile (5) — กระเบื้อง · ปูนกาว · ปูนยาแนว · คิ้ว", AMBER),
    ("column_beam (4) — ปูน · ทราย · หิน · น้ำ", TEAL),
    ("rebar (8) — RB6-RB9 · DB10-DB25 · ลวด · ไม้แบบ · ตะปู", RED),
]
for i, (txt, col) in enumerate(tax):
    y = Inches(6.3 + i * 0.3)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.07), Inches(0.13), Inches(0.13))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(12), Inches(0.3), txt, size=11.5, color=INK2)
add_footer(s, 3, TOTAL)

# ============================================================
# 4. ACQUISITION
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 04 · ACQUISITION")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
         "Four Strategies", size=32, bold=True, color=INK)

strategies = [
    ("A · PDF",       "unpdf parse",   "TPSO + CGD ส่งราคาใน PDF rate report. ใช้ unpdf + regex หาแถวที่ตรง material code",
     "Input: PDF URL\nOutput: number", AMBER),
    ("B · JSON API",  "suggest.jsp",   "HomePro/MegaHome expose /service/search/suggest.jsp?q=… JSON. ตรง, เร็ว, ไม่ต้อง render",
     "Input: search term\nOutput: {item_name, price}[]", TEAL),
    ("C · PROXY",     "ScrapingBee",   "Residential proxy + JS render. เก็บ JSON-LD Product → name+price แล้วกรองด้วย canonical",
     "Input: URL + waitMs\nOutput: HTML → ld+json", LINE),
    ("D · MANUAL",    "Admin upload",  "POST /api/admin/upload-prices รับ CSV/JSON ใส่ KV โดยตรง. ใช้กับ govt/SPA ที่ scrape ไม่ได้",
     "Input: {source,province,prices}\nOutput: KV write", RED),
]
sw = Inches(3.0); sh_ = Inches(3.0); sgap = Inches(0.1)
sx = Inches(0.6); sy = Inches(2.0)
for i, (tag, title, body, io, accent) in enumerate(strategies):
    x = sx + (sw + sgap) * i
    add_card(s, x, sy, sw, sh_, border=accent, border_pt=1.5)
    pill_color = AMBER_BR if accent is AMBER else (TEAL_LT if accent is TEAL else (RED_LT if accent is RED else INK2))
    pill_bg = RGBColor(0x2A,0x22,0x10) if accent is AMBER else (
              RGBColor(0x0E,0x2A,0x26) if accent is TEAL else (
              RGBColor(0x2A,0x10,0x10) if accent is RED else RGBColor(0x22,0x26,0x2E)))
    add_pill(s, x + Inches(0.2), sy + Inches(0.2), tag, color=pill_color, bg=pill_bg)
    add_text(s, x + Inches(0.2), sy + Inches(0.6), sw - Inches(0.4), Inches(0.4),
             title, size=15, bold=True, color=INK)
    add_text(s, x + Inches(0.2), sy + Inches(1.05), sw - Inches(0.4), Inches(1.3),
             body, size=10.5, color=INK2)
    add_text(s, x + Inches(0.2), sy + Inches(2.3), sw - Inches(0.4), Inches(0.7),
             io, size=9, color=INK3, font=FONT_MONO)

# Fallback chain
add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
         "Fallback chain (retail provider)", size=14, bold=True, color=INK)
fb = [
    ("1", "ScrapingBee + JSON-LD canonical pick", TEAL_LT),
    ("2", "ScrapingBee + regex extract",          AMBER_BR),
    ("3", "CF Browser Rendering (puppeteer)",     AMBER_BR),
    ("4", "null → wait for manual",               RED_LT),
]
fw = Inches(2.8); fh = Inches(1.0); fgap = Inches(0.18)
fx = Inches(0.6); fy = Inches(5.85)
for i, (n, t, col) in enumerate(fb):
    x = fx + (fw + fgap) * i
    add_card(s, x, fy, fw, fh)
    bg = RGBColor(0x0E,0x2A,0x26) if col is TEAL_LT else (
         RGBColor(0x2A,0x10,0x10) if col is RED_LT else RGBColor(0x2A,0x22,0x10))
    add_pill(s, x + Inches(0.15), fy + Inches(0.18), n, color=col, bg=bg)
    add_text(s, x + Inches(0.65), fy + Inches(0.18), fw - Inches(0.8), Inches(0.7),
             t, size=10, color=INK2)
    if i < len(fb) - 1:
        add_arrow(s, x + fw + Inches(0.005), fy + Inches(0.35), fgap - Inches(0.01), Inches(0.3))
add_footer(s, 4, TOTAL)

# ============================================================
# 5. CANONICALIZATION
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 05 · NORMALIZE")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
         "Option B — Canonical Spec Matching", size=28, bold=True, color=INK)
add_text(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.9),
         "ปัญหา: HomePro ค้นด้วยคำ \"ปูน\" ได้ผลลัพธ์ 200 รายการ ราคา 80-2,400฿. แต่ละแหล่งใช้คำต่างกัน "
         "ขนาดต่างกัน. ถ้า median ทั้งหมด = เปรียบเทียบไม่ตรง", size=14, color=INK2)

# Two cards
cw = Inches(6.0); ch = Inches(2.7)
add_card(s, Inches(0.6), Inches(2.95), cw, ch)
add_card(s, Inches(6.8), Inches(2.95), cw, ch)

add_text(s, Inches(0.85), Inches(3.1), Inches(5.5), Inches(0.4),
         "1. Canonical spec per material", size=15, bold=True, color=INK)
code = (
    "CEMENT_001: {\n"
    "  canonical: {\n"
    "    brand: \"ตราเสือ\",\n"
    "    size:  \"50kg\",\n"
    "    grade: \"Type I\"\n"
    "  },\n"
    "  searchTerms: [\n"
    "    \"ปูนซีเมนต์ปอร์ตแลนด์ ตราเสือ 50 กก\",\n"
    "    \"ปูนตราเสือ 50 กก\"\n"
    "  ]\n"
    "}"
)
add_card(s, Inches(0.85), Inches(3.55), Inches(5.5), Inches(2.0), fill=BG)
add_text(s, Inches(1.05), Inches(3.7), Inches(5.1), Inches(1.85),
         code, size=10, color=INK2, font=FONT_MONO)

add_text(s, Inches(7.05), Inches(3.1), Inches(5.5), Inches(0.4),
         "2. Three-tier filter", size=15, bold=True, color=INK)
tiers = [
    ("tight",  "ผลลัพธ์ที่ขนาดตรง (must) + brand/grade hit",  TEAL_LT, RGBColor(0x0E,0x2A,0x26)),
    ("loose",  "ขนาดตรงอย่างเดียว",                            AMBER_BR,RGBColor(0x2A,0x22,0x10)),
    ("legacy", "ไม่มี hit เลย → median ทั้งหมด (backwards compat)", INK2,  RGBColor(0x22,0x26,0x2E)),
]
for i, (p, t, c, bg) in enumerate(tiers):
    y = Inches(3.65 + i * 0.5)
    add_pill(s, Inches(7.05), y, p, color=c, bg=bg)
    add_text(s, Inches(8.05), y, Inches(4.5), Inches(0.4), t, size=11, color=INK2)
add_text(s, Inches(7.05), Inches(5.2), Inches(5.5), Inches(0.4),
         "ใช้ใน _suggestJsp.ts (HomePro/MegaHome) และ _scrapingbee.ts JSON-LD path",
         size=10, color=INK3, font=FONT_MONO)

# 3 honesty cards
add_text(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.4),
         "3. UI honesty", size=15, bold=True, color=INK)
honesty = [
    ("Canonical chip", "แสดง brand/size/grade ที่ใช้ filter", AMBER),
    ("Trust label",    "Government = Official BoQ · Modern Trade = Market", TEAL),
    ("Spread warning", "เมื่อ spread > 30% → \"Indicative range\"", RED),
]
hw = Inches(4.0); hh = Inches(0.95); hgap = Inches(0.1)
for i, (t, body, col) in enumerate(honesty):
    x = Inches(0.6) + (hw + hgap) * i
    y = Inches(6.3)
    add_card(s, x, y, hw, hh, border=col, border_pt=1.5)
    add_text(s, x + Inches(0.2), y + Inches(0.12), hw - Inches(0.4), Inches(0.35),
             t, size=12, bold=True, color=INK)
    add_text(s, x + Inches(0.2), y + Inches(0.5), hw - Inches(0.4), Inches(0.4),
             body, size=10.5, color=INK2)
add_footer(s, 5, TOTAL)

# ============================================================
# 6. STORAGE
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 06 · STORAGE")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
         "Cloudflare KV — Cache + Time-Series", size=30, bold=True, color=INK)

cw = Inches(6.0); ch = Inches(2.5)
add_card(s, Inches(0.6), Inches(2.0), cw, ch)
add_text(s, Inches(0.85), Inches(2.15), Inches(5.5), Inches(0.4),
         "Live price cache", size=15, bold=True, color=INK)
code1 = (
    "key: \"{source}:{materialId}:{province}\"\n"
    "val: { price: 175, fetchedAt: 1716... }\n"
    "TTL: 30 days (govt) / 14 days (HomePro)"
)
add_card(s, Inches(0.85), Inches(2.6), Inches(5.5), Inches(1.2), fill=BG)
add_text(s, Inches(1.05), Inches(2.75), Inches(5.1), Inches(1.05),
         code1, size=10.5, color=INK2, font=FONT_MONO)
add_text(s, Inches(0.85), Inches(3.95), Inches(5.5), Inches(0.5),
         "อ่านที่ /api/prices/[source]/[material] ทุก request → KV hit เป็นมาตรฐาน",
         size=10.5, color=INK3)

add_card(s, Inches(6.8), Inches(2.0), cw, ch)
add_text(s, Inches(7.05), Inches(2.15), Inches(5.5), Inches(0.4),
         "Weekly history snapshot", size=15, bold=True, color=INK)
code2 = (
    "key: \"hist:{source}:{material}:{prov}:{yyyy-mm-dd}\"\n"
    "val: { price, fetchedAt }\n"
    "TTL: 365 days"
)
add_card(s, Inches(7.05), Inches(2.6), Inches(5.5), Inches(1.2), fill=BG)
add_text(s, Inches(7.25), Inches(2.75), Inches(5.1), Inches(1.05),
         code2, size=10.5, color=INK2, font=FONT_MONO)
add_text(s, Inches(7.05), Inches(3.95), Inches(5.5), Inches(0.5),
         "เขียนทุกจันทร์ 03:17 UTC โดย GH Actions cron → ใช้สร้าง trend chart",
         size=10.5, color=INK3)

# Refresh cadence table
add_text(s, Inches(0.6), Inches(4.85), Inches(12), Inches(0.4),
         "Refresh cadence", size=15, bold=True, color=INK)
headers = ["TRIGGER", "WHAT", "FREQUENCY", "COST GUARD"]
cols_x = [Inches(0.6), Inches(3.6), Inches(7.2), Inches(9.6)]
cols_w = [Inches(2.95), Inches(3.55), Inches(2.35), Inches(3.6)]
for h, x, w in zip(headers, cols_x, cols_w):
    add_text(s, x, Inches(5.3), w, Inches(0.3), h, size=9, bold=True, color=AMBER_BR, font=FONT_MONO)
rows = [
    ("GH Actions cron", "refresh-prices + snapshot-history", "Mon 03:17 UTC", "~50 ScrapingBee credits/wk"),
    ("On-demand fetch", "/api/prices/[src]/[mat]",           "per request (KV miss only)", "rate-limited"),
    ("Admin upload",    "/api/admin/upload-prices",          "manual (CGD/DIT/SPAs)", "auth-gated"),
]
for i, row in enumerate(rows):
    y = Inches(5.7 + i * 0.42)
    sep = s.shapes.add_connector(1, Inches(0.6), y - Inches(0.04), Inches(13.2), y - Inches(0.04))
    sep.line.color.rgb = LINE; sep.line.width = Pt(0.5)
    for v, x, w in zip(row, cols_x, cols_w):
        add_text(s, x, y, w, Inches(0.4), v, size=10.5, color=INK2,
                 font=FONT_MONO if x == cols_x[0] else FONT_BODY)
add_footer(s, 6, TOTAL)

# ============================================================
# 7. API LAYER
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 07 · API LAYER")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
         "Endpoints", size=32, bold=True, color=INK)

# Two columns
col_w = Inches(6.0); col_h = Inches(4.3)
add_card(s, Inches(0.6), Inches(2.0), col_w, col_h)
add_card(s, Inches(6.8), Inches(2.0), col_w, col_h)

add_text(s, Inches(0.85), Inches(2.15), Inches(5.5), Inches(0.4),
         "Public REST", size=15, bold=True, color=TEAL_LT)
pub = [
    ("GET /api/prices/[source]/[material]?province=N", "→ {price, live, available, fetchedAt}"),
    ("GET /api/compare/[material]?province=N", "→ fan-out 10 sources + summary {min,max,spreadPct}"),
    ("GET /api/sources/health?province=N", "→ per-source {fresh,ok,stale,missing} + coverage%"),
    ("GET /api/prices/status", "→ live providers + ScrapingBee enabled?"),
    ("GET /api/trend/[source]/[material]", "→ time-series จาก hist: keys"),
]
for i, (ep, body) in enumerate(pub):
    y = Inches(2.65 + i * 0.7)
    add_text(s, Inches(0.85), y, Inches(5.5), Inches(0.3), ep, size=10, color=AMBER_BR, font=FONT_MONO)
    add_text(s, Inches(0.85), y + Inches(0.3), Inches(5.5), Inches(0.3), body, size=10, color=INK3)

add_text(s, Inches(7.05), Inches(2.15), Inches(5.5), Inches(0.4),
         "Admin (x-admin-token)", size=15, bold=True, color=AMBER_BR)
adm = [
    ("POST /api/admin/refresh-prices?source=X", "→ trigger live fetch + KV write"),
    ("POST /api/admin/upload-prices", "→ bulk ingest {source, province, prices{}}"),
    ("POST /api/admin/snapshot-history", "→ write hist: keys"),
    ("GET  /api/admin/scrapingbee-debug?source=X&material=Y", "→ surface upstream status, html signals, probes"),
]
for i, (ep, body) in enumerate(adm):
    y = Inches(2.65 + i * 0.7)
    add_text(s, Inches(7.05), y, Inches(5.5), Inches(0.3), ep, size=10, color=AMBER_BR, font=FONT_MONO)
    add_text(s, Inches(7.05), y + Inches(0.3), Inches(5.5), Inches(0.3), body, size=10, color=INK3)

# Bottom callout
add_card(s, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.7), border=TEAL, border_pt=1.5)
add_text(s, Inches(0.85), Inches(6.7), Inches(11.7), Inches(0.4),
         "Cache headers: public, s-maxage=60, stale-while-revalidate=120 — "
         "Cloudflare edge cache ตอบ 99% ของ traffic โดยไม่แตะ Worker",
         size=11, color=INK2)
add_footer(s, 7, TOTAL)

# ============================================================
# 8. CONSUMERS
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 08 · CONSUMERS")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
         "UI Pages", size=32, bold=True, color=INK)

pages = [
    ("/sources",        "ตารางราคา per material × per source + freshness badge + CSV export + manual uploader"),
    ("/compare-sources","เลือก material × province → bar chart + table 10 sources + canonical chip + spread warning"),
    ("/trend",          "Time-series line chart ดู price drift ของ material ต่อ source (รอ cron สะสม ≥7 จุด)"),
    ("/health",         "5 summary tiles + per-source coverage % (อ่าน KV ไม่ scrape — ฟรี)"),
    ("/api-docs",       "Public REST reference + curl examples + materials/sources enum"),
    ("/calc",           "BoQ calculator — เลือก work_type + ปริมาณ → multiply with cons → total"),
    ("/admin/sources",  "CsvUploader UI สำหรับ bulk ingest CGD/DIT/SPA monthly"),
    ("sitemap + robots","SEO. Disallow: /api/admin/"),
]
pw = Inches(3.0); ph = Inches(2.15); pgap = Inches(0.13)
for i, (tag, body) in enumerate(pages):
    col = i % 4
    row = i // 4
    x = Inches(0.6) + (pw + pgap) * col
    y = Inches(2.05) + (ph + pgap) * row
    add_card(s, x, y, pw, ph)
    add_pill(s, x + Inches(0.18), y + Inches(0.2), tag, color=AMBER_BR, bg=RGBColor(0x2A,0x22,0x10))
    add_text(s, x + Inches(0.2), y + Inches(0.7), pw - Inches(0.4), ph - Inches(0.8),
             body, size=10.5, color=INK2)
add_footer(s, 8, TOTAL)

# ============================================================
# 9. PROBLEMS
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 09 · PROBLEMS & MITIGATIONS")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
         "Known Gaps", size=32, bold=True, color=INK)

headers = ["PROBLEM", "ROOT CAUSE", "MITIGATION TODAY", "STATUS"]
cols_x = [Inches(0.6), Inches(4.3), Inches(7.7), Inches(11.5)]
cols_w = [Inches(3.6), Inches(3.3), Inches(3.7), Inches(1.7)]
for h, x, w in zip(headers, cols_x, cols_w):
    add_text(s, x, Inches(1.95), w, Inches(0.3), h, size=10, bold=True, color=AMBER_BR, font=FONT_MONO)
hbar = s.shapes.add_connector(1, Inches(0.6), Inches(2.3), Inches(13.2), Inches(2.3))
hbar.line.color.rgb = LINE; hbar.line.width = Pt(1)

rows = [
    ("CGD blocks scrape (403)",  "anti-bot",                    "Manual upload monthly",       ("WORKAROUND", RED)),
    ("DIT URL deprecated",       "infra change",                "Manual upload daily",         ("WORKAROUND", RED)),
    ("Thai Watsadu / BnB",       "Cloudflare bot challenge",    "Manual / paid residential proxy",("BLOCKED", RED)),
    ("Globalhouse/Dohome/SCG",   "SPA — initial HTML ว่าง, XHR หลัง hydrate", "Manual; รอ XHR reverse-engineering", ("BLOCKED", RED)),
    ("ScrapingBee free 1000/mo", "quota",                       "Weekly cron + canonical filter = ~200/mo", ("OK", TEAL)),
    ("Apples-to-apples matching","different brand/size per source","Option B canonical filter",  ("FIXED", TEAL)),
    ("Trend chart empty",        "cron เพิ่งเปลี่ยนเป็น weekly", "passive — สะสมไป 4-7 สัปดาห์",  ("WAITING", AMBER)),
]
for i, (p, rc, mt, (st, col)) in enumerate(rows):
    y = Inches(2.5 + i * 0.5)
    add_text(s, cols_x[0], y, cols_w[0], Inches(0.45), p, size=10.5, color=INK)
    add_text(s, cols_x[1], y, cols_w[1], Inches(0.45), rc, size=10, color=INK2)
    add_text(s, cols_x[2], y, cols_w[2], Inches(0.45), mt, size=10, color=INK2)
    bg = RGBColor(0x2A,0x10,0x10) if col is RED else (
         RGBColor(0x0E,0x2A,0x26) if col is TEAL else RGBColor(0x2A,0x22,0x10))
    pill_c = RED_LT if col is RED else (TEAL_LT if col is TEAL else AMBER_BR)
    add_pill(s, cols_x[3], y, st, color=pill_c, bg=bg)
    sep = s.shapes.add_connector(1, Inches(0.6), y + Inches(0.46), Inches(13.2), y + Inches(0.46))
    sep.line.color.rgb = LINE; sep.line.width = Pt(0.4)

add_card(s, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.85), border=AMBER, border_pt=1.5)
add_text(s, Inches(0.85), Inches(6.5), Inches(11.7), Inches(0.6),
         "Reality check: ~5/10 sources ขึ้น auto, 5/10 ต้อง manual. แต่ 3/10 ที่ auto "
         "(TPSO + HomePro + MegaHome) ครอบคลุม 70% ของ material lookups จริงในงาน BoQ",
         size=11.5, color=INK)
add_footer(s, 9, TOTAL)

# ============================================================
# 10. ROADMAP
# ============================================================
s = new_slide()
add_tag(s, Inches(0.6), Inches(0.5), "PIPELINE / 10 · ROADMAP")
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
         "Next Moves", size=32, bold=True, color=INK)

# Two columns
add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
         "Short-term (วันนี้ → 4 สัปดาห์)", size=15, bold=True, color=TEAL_LT)
short = [
    "1. รัน weekly cron 4 สัปดาห์ → trend page มีข้อมูลให้ดู",
    "2. เก็บ CGD/DIT รายไตรมาส โดย /admin/sources CSV uploader",
    "3. ขยาย canonical จาก 15/17 → 17/17 materials",
    "4. เขียน sourceOverrides ตัวอย่างใน 2-3 materials",
]
for i, t in enumerate(short):
    y = Inches(2.5 + i * 0.45)
    add_text(s, Inches(0.6), y, Inches(6), Inches(0.4), t, size=11.5, color=INK2)

add_text(s, Inches(6.8), Inches(2.0), Inches(6), Inches(0.4),
         "Medium-term (4-12 สัปดาห์)", size=15, bold=True, color=AMBER_BR)
mid = [
    "5. Reverse-engineer Dohome/SCG XHR API → unblock 2 sources",
    "6. Upgrade ScrapingBee → Hobby ($49/mo) ถ้า usage > 1000",
    "7. Province aggregator — median ราคา per region",
    "8. Alert: spread > 50% → push to Slack webhook",
]
for i, t in enumerate(mid):
    y = Inches(2.5 + i * 0.45)
    add_text(s, Inches(6.8), y, Inches(6), Inches(0.4), t, size=11.5, color=INK2)

# Open questions
add_text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.4),
         "Open questions", size=15, bold=True, color=INK)
oq = [
    ("Trust weighting?",   "Government = 2x weight, Retail = 1x?"),
    ("Province coverage?", "ตอนนี้ default 10 (กทม.) — ขยายกี่จังหวัด"),
    ("Public docs?",       "เปิด /api-docs สำหรับ 3rd-party + rate-limit policy"),
]
qw = Inches(4.0); qh = Inches(1.4); qgap = Inches(0.1)
for i, (t, body) in enumerate(oq):
    x = Inches(0.6) + (qw + qgap) * i
    y = Inches(5.15)
    add_card(s, x, y, qw, qh)
    add_text(s, x + Inches(0.2), y + Inches(0.2), qw - Inches(0.4), Inches(0.4),
             t, size=13, bold=True, color=INK)
    add_text(s, x + Inches(0.2), y + Inches(0.65), qw - Inches(0.4), qh - Inches(0.75),
             body, size=10.5, color=INK2)

# Footer line
sep = s.shapes.add_connector(1, Inches(0.6), Inches(6.85), Inches(13.2), Inches(6.85))
sep.line.color.rgb = LINE; sep.line.width = Pt(0.5)
add_text(s, Inches(0.6), Inches(6.95), Inches(13), Inches(0.4),
         "Live: construction-cost-engine.steep-tooth-c420.workers.dev · "
         "Repo: Bellero-Advanced/construction-cost-engine · Last: e8b76fc (Pipeline deck)",
         size=10, color=INK3, font=FONT_MONO)
add_footer(s, 10, TOTAL)

# ============================================================
out = "docs/pipeline-deck/Construction-Cost-Engine-Pipeline.pptx"
prs.save(out)
print(f"Wrote {out}")
